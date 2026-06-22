#!/usr/bin/env python3
"""Batch 2 of the blank-space evidence cards (native, editable). Each page
reproduces the slide's CURRENT text verbatim and adds one real-data card:

  page 1 -> slide 7   "How Trustworthiness Was Achieved" + guardrail rulebook
  page 2 -> slide 11  "The Discovery Challenge"          + invisible-inventory gap
  page 3 -> slide 14  "Defining Negative Space"          + Elks Run exclusion
  page 4 -> slide 16  "Shipping It Safely"               + four-layer key safety
  page 5 -> slide 18  "Roadmap: What I'd Do Next"        + impact-ranked ladder

Sources: agent/guardrails.py, agent/eval/backtest_report.md,
agent/eval/sample_prospect_interrupt.json, agent/demo_api.py.
Output: deck/assets/fill_batch2.pptx
"""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn

OUT = os.path.join(os.path.dirname(__file__), "assets"); os.makedirs(OUT, exist_ok=True)
BG="2B2A45"; PANEL="38365C"; PANEL2="454270"; PANEL_DK="232239"; TERM="1E1D33"
YEL="F0C84C"; YEL_LT="F7DD85"; LAV="A89FE6"; LAV_LT="C5BEF2"
GREEN="5FB98C"; RED="E0685C"; RED_LT="EBA59C"; AMBER="E2A24A"
FG="FFFFFF"; CODE="E8E7F2"; MUTE="C2C0D6"; FAINT="8E8CAB"; LINE="55527E"
HEAD="Poppins"; BODY="Poppins"; MONO="Roboto Mono"
PXEMU=6350
def E(px): return Emu(int(round(px*PXEMU)))
def RGB(h): return RGBColor.from_string(h)
CEN=PP_ALIGN.CENTER; RT=PP_ALIGN.RIGHT; MID=MSO_ANCHOR.MIDDLE; TOP=MSO_ANCHOR.TOP

prs=Presentation(); prs.slide_width=E(1920); prs.slide_height=E(1080); BLANK=prs.slide_layouts[6]
def slide():
    s=prs.slides.add_slide(BLANK)
    bg=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,E(0),E(0),E(1920),E(1080))
    bg.fill.solid(); bg.fill.fore_color.rgb=RGB(BG); bg.line.fill.background(); bg.shadow.inherit=False
    return s
def P(t,size,color,b=False,font=BODY,align=PP_ALIGN.LEFT,sb=0,sa=1,ls=1.05):
    return {"t":t,"size":size,"color":color,"b":b,"font":font,"align":align,"sb":sb,"sa":sa,"ls":ls}
def MP(runs,align=PP_ALIGN.LEFT,sb=0,sa=1,ls=1.1):
    return {"runs":runs,"align":align,"sb":sb,"sa":sa,"ls":ls,"t":runs[0]["t"],"size":runs[0]["size"],"color":runs[0]["color"]}
def _tf(tf,paras,anchor,ml=14,mr=10,mt=6,mb=6):
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=E(ml); tf.margin_right=E(mr); tf.margin_top=E(mt); tf.margin_bottom=E(mb)
    for i,p in enumerate(paras):
        par=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        par.alignment=p.get("align",PP_ALIGN.LEFT); par.line_spacing=p.get("ls",1.05)
        par.space_after=Pt(p.get("sa",1)); par.space_before=Pt(p.get("sb",0))
        for run in p.get("runs",[p]):
            r=par.add_run(); r.text=run["t"]; f=r.font
            f.size=Pt(run["size"]); f.bold=run.get("b",False); f.name=run.get("font",BODY); f.color.rgb=RGB(run["color"])
def box(s,x,y,w,h,fill=PANEL,line=LINE,lw=1.6,radius=0.08,paras=None,anchor=MID,ml=14,mr=10,mt=6,mb=6,dash=False):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,E(x),E(y),E(w),E(h))
    try: sp.adjustments[0]=radius
    except Exception: pass
    sp.fill.solid(); sp.fill.fore_color.rgb=RGB(fill)
    if line:
        sp.line.color.rgb=RGB(line); sp.line.width=Pt(lw)
        if dash: sp.line._get_or_add_ln().append(sp.line._get_or_add_ln().makeelement(qn("a:prstDash"),{"val":"dash"}))
    else: sp.line.fill.background()
    sp.shadow.inherit=False
    if paras: _tf(sp.text_frame,paras,anchor,ml=ml,mr=mr,mt=mt,mb=mb)
    return sp
def label(s,x,y,w,h,paras,anchor=TOP):
    tb=s.shapes.add_textbox(E(x),E(y),E(w),E(h)); _tf(tb.text_frame,paras,anchor,ml=2,mr=2,mt=2,mb=2); return tb
def conn(s,x1,y1,x2,y2,color=LINE,width=1.2,arrow=False):
    c=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,E(x1),E(y1),E(x2),E(y2))
    c.line.color.rgb=RGB(color); c.line.width=Pt(width)
    if arrow: c.line._get_or_add_ln().append(c.line._get_or_add_ln().makeelement(qn("a:tailEnd"),{"type":"triangle","w":"med","len":"med"}))
    return c

def footer(s): label(s,70,1016,400,28,[P("BlueLiner Agents",11,FAINT,font=HEAD)])
def text_chrome(s,title,subtitle,body_paras,title_size=50,subtitle_top=300,title_h=230):
    label(s,70,66,884,title_h,[P(title,title_size,FG,b=True,font=HEAD,ls=1.0)])
    label(s,70,subtitle_top,770,60,[P(subtitle,23,FG,font=HEAD)])
    label(s,976,70,874,920,body_paras); footer(s)
def bod(t): return P(t,16.5,FG,ls=1.28,sa=12)
def lead(l,r,sa=12,size=16.5): return MP([{"t":l,"size":size,"color":FG,"b":True},{"t":r,"size":size,"color":FG}],ls=1.28,sa=sa)
def badge(s,x,y,n,color=YEL,d=30):
    box(s,x,y,d,d,fill=color,line=None,radius=0.5,anchor=MID,paras=[P(str(n),14,PANEL_DK,b=True,align=CEN,font=HEAD)])

# ===================== CARD: guardrail rulebook (slide 7) =====================
def guardrail_card(s,ox,oy,W=812,H=548):
    A=lambda x:ox+x; B=lambda y:oy+y
    box(s,A(0),B(0),W,H,fill=PANEL,line=LINE,lw=1.5,radius=0.03)
    label(s,A(24),B(14),W-40,20,[P("GROUNDING + GUARDRAILS",11.5,FAINT,b=True,font=HEAD)])
    box(s,A(24),B(42),W-48,66,fill=PANEL_DK,line=YEL,lw=1.4,radius=0.06,anchor=MID,ml=18,
        paras=[P("Grounding contract",13.5,YEL_LT,b=True,font=HEAD),
               MP([{"t":"every number traces to a tool result, else regenerate then strip   ","size":11.5,"color":MUTE},
                   {"t":"hallucinations 100% → 0%","size":11.5,"color":GREEN,"b":True}],sb=2)])
    label(s,A(24),B(120),W-40,20,[P("HARD GUARDRAILS — deterministic, non-overridable",11,FAINT,b=True,font=HEAD)])
    rules=[("Flood","flow > 3× median","BLOCK",RED),
           ("Too warm","water > 68°F","BLOCK",RED),
           ("Too cold","water < 40°F","DEMOTE",AMBER),
           ("Private access","no public entry","BLOCK",RED),
           ("Staleness","reading too old","DEMOTE",AMBER)]
    y=150
    for name,trig,act,col in rules:
        box(s,A(24),B(y),W-48,48,fill=PANEL_DK,line=col,lw=1.3,radius=0.1)
        label(s,A(44),B(y+13),260,24,[P(name,14,FG,b=True,font=HEAD)])
        label(s,A(300),B(y+14),340,24,[P(trig,12.5,MUTE,font=MONO)])
        label(s,A(W-180),B(y+13),150,24,[P(act,12.5,col,b=True,align=RT,font=HEAD)])
        y+=56
    label(s,A(24),B(H-36),W-48,28,[P("the model advises  ·  the rules decide",12,FAINT,align=CEN,font=HEAD)])

def slide7(s):
    text_chrome(s,"How Trustworthiness Was Achieved","Grounding and hard guardrails",[
        lead("Two mechanisms:"," grounding and hard guardrails."),
        lead("Grounding contract:"," every number must trace to a tool result. If it can't, regenerate once, then strip it. (Hallucinated readings went from 100% to 0%.)"),
        lead("Guardrails are deterministic code, not prompts:"," flood (flow over 3x the median), trout-ethics temperature band, private-access block, and staleness demotion."),
        lead("The model advises; the rules decide."," v3 cannot recommend blocked water by construction.")],
        title_size=44,subtitle_top=345,title_h=270)
    guardrail_card(s,70,432)

# ===================== CARD: invisible inventory (slide 11) =====================
def inventory_card(s,ox,oy,W=812,H=576):
    A=lambda x:ox+x; B=lambda y:oy+y
    box(s,A(0),B(0),W,H,fill=PANEL,line=LINE,lw=1.5,radius=0.03)
    label(s,A(24),B(16),W-40,20,[P("INVISIBLE INVENTORY",11.5,FAINT,b=True,font=HEAD)])
    label(s,A(24),B(54),W-40,20,[P("Designated — shown on the official maps",12,MUTE)])
    box(s,A(24),B(80),215,40,fill=GREEN,line=None,radius=0.25)
    label(s,A(250),B(88),300,26,[P("36,378 reaches",15,GREEN,b=True,font=HEAD)])
    label(s,A(24),B(140),W-40,20,[P("Undesignated — scanned by the Prospector",12,MUTE)])
    box(s,A(24),B(166),600,40,fill=YEL,line=None,radius=0.25)
    label(s,A(634),B(174),160,26,[P("100K+",15,YEL_LT,b=True,font=HEAD)])
    box(s,A(24),B(238),W-48,86,fill=PANEL_DK,line=YEL,lw=1.4,radius=0.06,anchor=MID,ml=18,
        paras=[P("The value is in qualifying inventory the market has mispriced.",14.5,FG,b=True,ls=1.25),
               P("the Red Ventures parallel",11.5,YEL_LT,sb=3)])
    conn(s,A(24),B(348),A(W-24),B(348),color=LINE)
    label(s,A(24),B(364),W-48,24,[P("Official maps list only state-designated trout water; plenty of fishable trout",11.5,FAINT)])
    label(s,A(24),B(386),W-48,24,[P("water is undesignated. The Prospector ranks that hidden pool.",11.5,FAINT)])

def slide11(s):
    text_chrome(s,"The Discovery Challenge","Finding water the maps miss",[
        bod("Discovery is finding undervalued inventory."),
        bod("Official maps show only designated trout water, but plenty of fishable trout water is undesignated. That's invisible inventory."),
        bod("Business analogy (Red Ventures): the value is in qualifying inventory the market has mispriced."),
        bod("The Prospector surfaces undesignated-but-fishable reaches, ranked and qualified.")])
    inventory_card(s,70,400)

# ===================== CARD: Elks Run exclusion (slide 14) =====================
def exclusion_card(s,ox,oy,W=588,H=648):
    A=lambda x:ox+x; B=lambda y:oy+y
    box(s,A(0),B(0),W,H,fill=PANEL,line=LINE,lw=1.5,radius=0.03)
    box(s,A(6),B(22),6,H-44,fill=RED,line=None,radius=0.5)
    label(s,A(28),B(14),W-50,20,[P("EXCLUDED — OBVIOUS, NOT USEFUL",11.5,FAINT,b=True,font=HEAD)])
    box(s,A(24),B(46),W-48,128,fill=PANEL_DK,line=RED,lw=1.4,radius=0.05,anchor=TOP,ml=18,mt=14,
        paras=[MP([{"t":"✗  ","size":16,"color":RED,"b":True},{"t":"Elks Run","size":16,"color":RED_LT,"b":True,"font":HEAD}]),
               P("undesignated reach on a stream already on the map",12.5,MUTE,sb=4),
               P("topology distance 0.0 mi — it is the map, relabeled",12,FAINT,sb=2)])
    label(s,A(24),B(192),W-48,20,[P("THE RULE",11,FAINT,b=True,font=HEAD)])
    box(s,A(24),B(216),W-48,70,fill=PANEL_DK,line=LINE,lw=1,radius=0.06,anchor=MID,ml=18,
        paras=[P("Exclude same-stream extensions — don't just relabel them.",13.5,FG,ls=1.25)])
    label(s,A(24),B(300),W-48,20,[P("RECEIPTS",11,FAINT,b=True,font=HEAD)])
    rec=[("−29%","candidates removed"),("0","distance-0 results (was many)"),("✓","real tributary leads surfaced")]
    y=326
    for big,lab in rec:
        box(s,A(24),B(y),W-48,60,fill=PANEL_DK,line=GREEN,lw=1.2,radius=0.1,anchor=MID,ml=18,
            paras=[MP([{"t":big+"   ","size":15,"color":GREEN,"b":True,"font":HEAD},
                       {"t":lab,"size":12.5,"color":MUTE}])])
        y+=70
    label(s,A(24),B(H-44),W-48,30,[P("obvious is as disqualifying as wrong",11.5,FAINT,align=CEN)])

def slide14(s):
    label(s,70,60,1340,120,[P("Defining Negative Space",48,FG,b=True,font=HEAD,ls=1.0)])
    label(s,70,205,1340,56,[P("Deciding what NOT to surface is half of product quality",22,FG,font=HEAD)])
    label(s,70,330,874,52,[P("What to Exclude",21,FG,b=True,font=HEAD)])
    body=[MP([{"t":"The bad result: ","size":15.5,"color":FG,"b":True},
              {"t":"Elks Run, an undesignated reach on a stream we already show on the map (distance 0.0 mi).","size":15.5,"color":MUTE}],ls=1.3,sa=11),
          MP([{"t":"The judgment: ","size":15.5,"color":FG,"b":True},
              {"t":"that's not a discovery. Obvious is as disqualifying as wrong.","size":15.5,"color":MUTE}],ls=1.3,sa=11),
          MP([{"t":"The decision: ","size":15.5,"color":FG,"b":True},
              {"t":"exclude same-stream extensions, don't just relabel them. A clearer label on a useless result is still useless.","size":15.5,"color":MUTE}],ls=1.3,sa=11),
          MP([{"t":"Receipts: ","size":15.5,"color":FG,"b":True},
              {"t":"we removed 29% of candidates, dropped distance-0 results to zero, and surfaced the real tributary leads. I reported the metric move with its cause.","size":15.5,"color":MUTE}],ls=1.3,sa=11),
          P("A recommender's credibility dies the first time it tells you something you already know.",15.5,YEL_LT,b=True,ls=1.3,sb=4)]
    label(s,70,420,874,560,body)
    conn(s,1580,60,1580,1020,color=LINE,width=1.0)
    footer(s)
    exclusion_card(s,966,330)

# ===================== CARD: four-layer key safety (slide 16) =====================
def layers_card(s,ox,oy,W=812,H=576):
    A=lambda x:ox+x; B=lambda y:oy+y
    box(s,A(0),B(0),W,H,fill=PANEL,line=LINE,lw=1.5,radius=0.03)
    label(s,A(24),B(16),W-40,20,[P("API-KEY SAFETY — FOUR INDEPENDENT LAYERS",11.5,FAINT,b=True,font=HEAD)])
    label(s,A(24),B(40),W-40,20,[P("any one layer is sufficient",12,YEL_LT)])
    layers=["The public app never mounts the endpoint",
            "Agent deps aren't in the production image",
            "No key on the web service",
            "Off by default — behind a flag + optional token"]
    y=74
    for i,t in enumerate(layers,1):
        box(s,A(24),B(y),W-48,68,fill=PANEL_DK,line=LINE,lw=1.2,radius=0.08)
        badge(s,A(42),B(y+19),i,color=YEL,d=30)
        label(s,A(92),B(y+21),W-150,26,[P(t,14,FG,font=HEAD)])
        y+=82
    box(s,A(24),B(y+6),W-48,66,fill=PANEL_DK,line=GREEN,lw=1.4,radius=0.06,anchor=MID,ml=18,
        paras=[MP([{"t":"Blast-radius insurance: ","size":13,"color":GREEN,"b":True},
                   {"t":"a dedicated, spend-capped, revocable key.","size":13,"color":FG}])])

def slide16(s):
    text_chrome(s,"Shipping It Safely","Bound the downside before it bites",[
        lead("The threat:"," the public app becoming a free, unmetered proxy to my API key.",sa=12),
        P("Four independent layers, any one sufficient:",16.5,FG,ls=1.28,sa=4),
        P("The public app never mounts the endpoint",15,MUTE,ls=1.2,sa=2),
        P("Agent deps aren't in the production image",15,MUTE,ls=1.2,sa=2),
        P("No key on the web service",15,MUTE,ls=1.2,sa=2),
        P("Off by default, behind a flag and optional token",15,MUTE,ls=1.2,sa=12),
        lead("Blast-radius insurance:"," a dedicated, spend-capped, revocable key.",sa=12),
        bod("Can I 110% guarantee it's safe? No one can, so I bounded the downside instead.")])
    layers_card(s,70,400)

# ===================== CARD: roadmap ladder (slide 18) =====================
def roadmap_card(s,ox,oy,W=812,H=576):
    A=lambda x:ox+x; B=lambda y:oy+y
    box(s,A(0),B(0),W,H,fill=PANEL,line=LINE,lw=1.5,radius=0.03)
    label(s,A(24),B(16),W-40,20,[P("WHAT I'D DO NEXT",11.5,FAINT,b=True,font=HEAD)])
    rungs=[("Human-in-the-loop in the UI","make the headless confirm interactive (interrupt + resume)",LINE,None),
           ("Close the access data gap","wire PAD-US public-land polygons",YEL,"highest leverage · the binding constraint"),
           ("Exact flow-network topology (NLDI)","for the shortlist ranking",LINE,None),
           ("Grow the confirm/deny flywheel","calibration improves as anglers confirm",LINE,None)]
    y=52; n=1
    for title,sub,col,tag in rungs:
        h=82
        box(s,A(24),B(y),W-48,h,fill=PANEL_DK,line=col,lw=(2.2 if col==YEL else 1.2),radius=0.07)
        badge(s,A(42),B(y+26),n,color=(YEL if col==YEL else LINE),d=30)
        label(s,A(92),B(y+14),W-150,26,[P(title,14.5,(YEL_LT if col==YEL else FG),b=True,font=HEAD)])
        label(s,A(92),B(y+44),W-150,24,[P(sub,12,MUTE)])
        if tag: label(s,A(92),B(y+62),W-150,20,[P(tag,11,YEL_LT,b=True)])
        y+=h+12; n+=1
    label(s,A(24),B(H-40),W-48,28,[P("the eval pointed at the roadmap, not at a better model",12.5,YEL_LT,b=True,align=CEN,font=HEAD)])

def slide18(s):
    text_chrome(s,"Roadmap: What I'd Do Next","Highest-leverage moves first",[
        bod("Human in the Loop integrated into UI/UX"),
        bod("Close the access data gap by wiring public-land polygons. The bottleneck is data coverage of the binding constraint, not modeling."),
        bod("Use exact flow-network topology (NLDI) for the shortlist."),
        bod("Grow the confirm/deny flywheel so calibration improves as anglers confirm."),
        bod("The eval pointed at the roadmap, not at a better model.")])
    roadmap_card(s,70,400)

slide7(slide()); slide11(slide()); slide14(slide()); slide16(slide()); slide18(slide())
path=os.path.join(OUT,"fill_batch2.pptx"); prs.save(path)
print("wrote",path,os.path.getsize(path),"bytes,",len(prs.slides._sldIdLst),"slides")
