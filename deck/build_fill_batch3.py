#!/usr/bin/env python3
"""Batch 3: a new Context & Contracts slide + How Trustworthiness Was Achieved
rebuilt with production monitoring folded in. Native, editable (same pipeline).

  page 1 -> NEW "Context & Contracts" (insert after Two Agents / p4)
  page 2 -> rebuilt "How Trustworthiness Was Achieved" (+ monitoring) -> swaps p7

Sources: agent/agent.py (REC_SCHEMA, prompts), agent/guardrails.py, agent/memory.py,
agent/observability.py (RunTrace), agent/watch.py (proactive alerts), demo_api.py.
Output: deck/assets/fill_batch3.pptx
"""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn

OUT = os.path.join(os.path.dirname(__file__), "assets"); os.makedirs(OUT, exist_ok=True)
BG="2B2A45"; PANEL="38365C"; PANEL2="454270"; PANEL_DK="232239"; TERM="1E1D33"
YEL="F0C84C"; YEL_LT="F7DD85"; LAV="A89FE6"; LAV_LT="C5BEF2"
GREEN="5FB98C"; RED="E0685C"; RED_LT="EBA59C"; AMBER="E2A24A"
FG="FFFFFF"; CODE="E8E7F2"; MUTE="C2C0D6"; FAINT="8E8CAB"; LINE="55527E"
HEAD="Poppins"; BODY="Poppins"; MONO="Roboto Mono"
PXEMU=6350
def E(px): return Emu(int(round(px*PXEMU)))
def RGB(h): return RGBColor.from_string(h)
CEN=PP_ALIGN.CENTER; RT=PP_ALIGN.RIGHT; MID=MSO_ANCHOR.MIDDLE; TOP=MSO_ANCHOR.TOP

prs=Presentation(); prs.slide_width=E(1920); prs.slide_height=E(1080); BLANK=prs.slide_layouts[6]
def slide():
    s=prs.slides.add_slide(BLANK)
    bg=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,E(0),E(0),E(1920),E(1080))
    bg.fill.solid(); bg.fill.fore_color.rgb=RGB(BG); bg.line.fill.background(); bg.shadow.inherit=False
    return s
def P(t,size,color,b=False,font=BODY,align=PP_ALIGN.LEFT,sb=0,sa=1,ls=1.05):
    return {"t":t,"size":size,"color":color,"b":b,"font":font,"align":align,"sb":sb,"sa":sa,"ls":ls}
def MP(runs,align=PP_ALIGN.LEFT,sb=0,sa=1,ls=1.1):
    return {"runs":runs,"align":align,"sb":sb,"sa":sa,"ls":ls,"t":runs[0]["t"],"size":runs[0]["size"],"color":runs[0]["color"]}
def _tf(tf,paras,anchor,ml=14,mr=10,mt=6,mb=6):
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=E(ml); tf.margin_right=E(mr); tf.margin_top=E(mt); tf.margin_bottom=E(mb)
    for i,p in enumerate(paras):
        par=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        par.alignment=p.get("align",PP_ALIGN.LEFT); par.line_spacing=p.get("ls",1.05)
        par.space_after=Pt(p.get("sa",1)); par.space_before=Pt(p.get("sb",0))
        for run in p.get("runs",[p]):
            r=par.add_run(); r.text=run["t"]; f=r.font
            f.size=Pt(run["size"]); f.bold=run.get("b",False); f.name=run.get("font",BODY); f.color.rgb=RGB(run["color"])
def box(s,x,y,w,h,fill=PANEL,line=LINE,lw=1.6,radius=0.08,paras=None,anchor=MID,ml=14,mr=10,mt=6,mb=6):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,E(x),E(y),E(w),E(h))
    try: sp.adjustments[0]=radius
    except Exception: pass
    sp.fill.solid(); sp.fill.fore_color.rgb=RGB(fill)
    if line: sp.line.color.rgb=RGB(line); sp.line.width=Pt(lw)
    else: sp.line.fill.background()
    sp.shadow.inherit=False
    if paras: _tf(sp.text_frame,paras,anchor,ml=ml,mr=mr,mt=mt,mb=mb)
    return sp
def label(s,x,y,w,h,paras,anchor=TOP):
    tb=s.shapes.add_textbox(E(x),E(y),E(w),E(h)); _tf(tb.text_frame,paras,anchor,ml=2,mr=2,mt=2,mb=2); return tb
def conn(s,x1,y1,x2,y2,color=LINE,width=1.2):
    c=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,E(x1),E(y1),E(x2),E(y2))
    c.line.color.rgb=RGB(color); c.line.width=Pt(width); return c
def footer(s): label(s,70,1016,400,28,[P("BlueLiner Agents",11,FAINT,font=HEAD)])
def text_chrome(s,title,subtitle,body,title_size=50,subtitle_top=300,title_h=230):
    label(s,70,66,884,title_h,[P(title,title_size,FG,b=True,font=HEAD,ls=1.0)])
    label(s,70,subtitle_top,800,60,[P(subtitle,23,FG,font=HEAD)])
    label(s,976,70,874,930,body); footer(s)
def lead(l,r,size=16,sa=11): return MP([{"t":l,"size":size,"color":FG,"b":True},{"t":r,"size":size,"color":FG}],ls=1.28,sa=sa)
def bod(t,size=16): return P(t,size,FG,ls=1.28,sa=11)

# ================= NEW SLIDE: Context & Contracts =================
def context_card(s,ox,oy,W=812,H=576):
    A=lambda x:ox+x; B=lambda y:oy+y
    box(s,A(0),B(0),W,H,fill=PANEL,line=LINE,lw=1.5,radius=0.03)
    label(s,A(22),B(14),W-40,20,[P("WHAT THE AGENT NEEDS TO KNOW",11.5,FAINT,b=True,font=HEAD)])
    cw=377; gap=14; x1=22; x2=22+cw+gap; y1=46; ch=247; y2=y1+ch+14
    def cell(cx,cy,head,items):
        box(s,A(cx),B(cy),cw,ch,fill=PANEL_DK,line=LINE,lw=1.2,radius=0.05,anchor=TOP,ml=16,mt=14,mr=12,
            paras=[P(head,12.5,YEL_LT,b=True,font=HEAD,sa=6)]+items)
    cell(x1,y1,"SCHEMAS · CONTRACTS",[
        P("REC_SCHEMA: verdict · score ·",11,FG,font=MONO,ls=1.2),
        P("confidence · why[] · sources[]",11,FG,font=MONO,ls=1.2,sa=5),
        P("typed MCP tool I/O",11.5,MUTE,sb=2),
        P("grounding: every number cites a source",11.5,MUTE,sb=4,ls=1.2)])
    cell(x2,y1,"PROMPTS · versioned in git",[
        P("system.md",11.5,FG,font=MONO,sa=2),
        P("ranker.md",11.5,FG,font=MONO,sa=2),
        P("prospector_system.md",11.5,FG,font=MONO,sa=6),
        P("files I can diff, not notebook strings",11.5,MUTE,ls=1.2)])
    cell(x1,y2,"DURABLE CONTEXT",[
        P("catch-log memory: your temp + flow bands",11.5,FG,ls=1.2,sa=5),
        P("30-yr USGS medians for today's date",11.5,MUTE,ls=1.2,sa=5),
        P("scorer domain constants (encoded knowledge)",11.5,MUTE,ls=1.2)])
    cell(x2,y2,"BUSINESS RULES · PERMISSIONS",[
        P("legality · ethics · flood · staleness",11.5,FG,ls=1.2,sa=5),
        P("deterministic, non-overridable",11.5,MUTE,sb=1,sa=5),
        P("gated, off-by-default access",11.5,MUTE)])

def slide_context(s):
    text_chrome(s,"Context & Contracts","What the agent needs to know to work",[
        bod("A trustworthy agent needs durable context and hard contracts, not just a capable model."),
        lead("Contracts:"," a structured output schema (every recommendation carries a verdict, score, confidence, reasons, and sources) plus typed MCP tool I/O. The grounding contract makes sources mandatory."),
        lead("Durable context:"," per-user catch-log memory (your proven temperature and flow bands), 30-year medians for today's date, and the scorer's domain constants. Thresholds are encoded knowledge, not guesses."),
        lead("Prompts as code:"," system, ranker, and prospector prompts are versioned files I can diff, not strings in a notebook."),
        lead("Rules and permissions:"," legality, ethics, flood, and staleness are deterministic; the agent runs under a gated, off-by-default permission model.")])
    context_card(s,70,400)

# ============ REBUILD: How Trustworthiness Was Achieved (+ monitoring) ============
def guardrail_monitor_card(s,ox,oy,W=812,H=470):
    A=lambda x:ox+x; B=lambda y:oy+y
    box(s,A(0),B(0),W,H,fill=PANEL,line=LINE,lw=1.5,radius=0.03)
    label(s,A(24),B(10),W-40,20,[P("GROUNDING + GUARDRAILS",11.5,FAINT,b=True,font=HEAD)])
    box(s,A(24),B(34),W-48,58,fill=PANEL_DK,line=YEL,lw=1.4,radius=0.06,anchor=MID,ml=18,
        paras=[P("Grounding contract",13,YEL_LT,b=True,font=HEAD),
               MP([{"t":"every number traces to a tool result, else regenerate then strip   ","size":11,"color":MUTE},
                   {"t":"hallucinations 100% → 0%","size":11,"color":GREEN,"b":True}],sb=2)])
    label(s,A(24),B(100),W-40,18,[P("HARD GUARDRAILS — deterministic, non-overridable",11,FAINT,b=True,font=HEAD)])
    rules=[("Flood","flow > 3× median","BLOCK",RED),("Too warm","water > 68°F","BLOCK",RED),
           ("Too cold","water < 40°F","DEMOTE",AMBER),("Private access","no public entry","BLOCK",RED),
           ("Staleness","reading too old","DEMOTE",AMBER)]
    y=124
    for name,trig,act,col in rules:
        box(s,A(24),B(y),W-48,38,fill=PANEL_DK,line=col,lw=1.2,radius=0.1)
        label(s,A(42),B(y+9),250,22,[P(name,13,FG,b=True,font=HEAD)])
        label(s,A(298),B(y+10),330,22,[P(trig,12,MUTE,font=MONO)])
        label(s,A(W-176),B(y+9),148,22,[P(act,12,col,b=True,align=RT,font=HEAD)])
        y+=44
    conn(s,A(24),B(350),A(W-24),B(350),color=LINE)
    box(s,A(24),B(362),W-48,86,fill=PANEL_DK,line=LAV,lw=1.3,radius=0.06,anchor=TOP,ml=18,mt=11,
        paras=[P("MONITORED IN PRODUCTION",10.5,LAV_LT,b=True,font=HEAD,sa=4),
               P("every veto logged with its reason · every decision traced (tools · latency · cost · grounding) · proactive condition alerts",11,MUTE,ls=1.25)])

def slide_trust(s):
    text_chrome(s,"How Trustworthiness Was Achieved","Grounding, guardrails, and live monitoring",[
        lead("Two mechanisms:"," grounding and hard guardrails."),
        lead("Grounding contract:"," every number must trace to a tool result. If it can't, regenerate once, then strip it. (Hallucinated readings went from 100% to 0%.)"),
        lead("Guardrails are deterministic code, not prompts:"," flood (flow over 3x the median), trout-ethics temperature band, private-access block, and staleness demotion."),
        lead("The model advises; the rules decide."," v3 cannot recommend blocked water by construction."),
        lead("And we watch it in production:"," every veto is logged with its reason, every decision is traced (tools, latency, cost, grounding), and a proactive watch flags condition changes.")],
        title_size=44,subtitle_top=345,title_h=270)
    guardrail_monitor_card(s,70,420)

slide_context(slide()); slide_trust(slide())
path=os.path.join(OUT,"fill_batch3.pptx"); prs.save(path)
print("wrote",path,os.path.getsize(path),"bytes,",len(prs.slides._sldIdLst),"slides")
