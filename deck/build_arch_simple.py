#!/usr/bin/env python3
"""Simplified architecture diagram — same structure, fewer annotations,
readable by a panel in ~20s. Native editable PowerPoint shapes -> import to Canva.
Output: deck/assets/architecture_simple.pptx (1 slide)
"""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn

OUT = os.path.join(os.path.dirname(__file__), "assets"); os.makedirs(OUT, exist_ok=True)
NAVY="0B2A3A"; PANEL="102E41"; PANEL2="143A50"; PANEL_DK="0C2535"
BLUE="5BA8C8"; BLUE_LT="95C5D9"; GREEN="4A8C5C"; GREEN_LT="7FBE8E"
OCHRE="B7892F"; CLAY="B3473B"; CLAY_LT="E0A59B"; PURPLE="7A3DB8"; PURPLE_LT="A878DA"
FG="EAF2F6"; MUTE="9DB6C2"; FAINT="6E8C9C"; LINE="3E5C6E"; HEAD="Poppins"; BODY="Poppins"; MONO="Roboto Mono"
PXEMU=6350
def E(px): return Emu(int(round(px*PXEMU)))
def RGB(h): return RGBColor.from_string(h)
prs=Presentation(); prs.slide_width=E(1920); prs.slide_height=E(1080)
s=prs.slides.add_slide(prs.slide_layouts[6])
bg=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,E(0),E(0),E(1920),E(1080)); bg.fill.solid(); bg.fill.fore_color.rgb=RGB(NAVY); bg.line.fill.background(); bg.shadow.inherit=False

def P(t,size,color,b=False,font=BODY,align=PP_ALIGN.LEFT,sb=0,sa=1,ls=1.05):
    return {"t":t,"size":size,"color":color,"b":b,"font":font,"align":align,"sb":sb,"sa":sa,"ls":ls}
def _tf(tf,paras,anchor,ml=14,mr=10,mt=6,mb=6):
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=E(ml); tf.margin_right=E(mr); tf.margin_top=E(mt); tf.margin_bottom=E(mb)
    for i,p in enumerate(paras):
        par=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        par.alignment=p.get("align",PP_ALIGN.LEFT); par.line_spacing=p.get("ls",1.05)
        par.space_after=Pt(p.get("sa",1)); par.space_before=Pt(p.get("sb",0))
        r=par.add_run(); r.text=p["t"]; f=r.font
        f.size=Pt(p["size"]); f.bold=p.get("b",False); f.name=p.get("font",BODY); f.color.rgb=RGB(p["color"])
def box(x,y,w,h,fill=PANEL,line=LINE,lw=1.6,radius=0.08,paras=None,anchor=MSO_ANCHOR.MIDDLE):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,E(x),E(y),E(w),E(h))
    try: sp.adjustments[0]=radius
    except Exception: pass
    sp.fill.solid(); sp.fill.fore_color.rgb=RGB(fill)
    if line: sp.line.color.rgb=RGB(line); sp.line.width=Pt(lw)
    else: sp.line.fill.background()
    sp.shadow.inherit=False
    if paras: _tf(sp.text_frame,paras,anchor)
    return sp
def label(x,y,w,h,paras,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(E(x),E(y),E(w),E(h)); _tf(tb.text_frame,paras,anchor,ml=2,mr=2,mt=2,mb=2); return tb
def conn(x1,y1,x2,y2,color=MUTE,width=2.0,dashed=False,arrow=True):
    c=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,E(x1),E(y1),E(x2),E(y2))
    c.line.color.rgb=RGB(color); c.line.width=Pt(width); ln=c.line._get_or_add_ln()
    if dashed: ln.append(ln.makeelement(qn("a:prstDash"),{"val":"dash"}))
    if arrow: ln.append(ln.makeelement(qn("a:tailEnd"),{"type":"triangle","w":"med","len":"med"}))
    return c
def hdr(cx,t): label(cx-200,150,400,28,[P(t,11,FAINT,b=True,align=PP_ALIGN.CENTER,font=HEAD)])

# title
label(70,34,1780,60,[P("Two Agents, One Trustworthy Spine",32,FG,b=True,font=HEAD)])
label(70,92,1780,30,[P("Two agents share one deterministic, guardrailed spine.",15,MUTE)])
hdr(185,"REQUESTS"); hdr(550,"ORCHESTRATION"); hdr(1015,"SHARED TRUSTWORTHY SPINE"); hdr(1560,"DATA SOURCES")

# Requests
box(70,210,230,120, line=BLUE, lw=2, paras=[P("Trip Planner",15,BLUE_LT,b=True,font=HEAD,align=PP_ALIGN.CENTER),P("“where to fish?”",11,MUTE,align=PP_ALIGN.CENTER,sb=2)])
box(70,470,230,120, line=PURPLE, lw=2, paras=[P("Prospector",15,PURPLE_LT,b=True,font=HEAD,align=PP_ALIGN.CENTER),P("“find new trout water?”",11,MUTE,align=PP_ALIGN.CENTER,sb=2)])

# Orchestration
box(360,210,420,150, line=BLUE, lw=2, anchor=MSO_ANCHOR.MIDDLE, paras=[
    P("Trip Planner",18,BLUE_LT,b=True,font=HEAD),P("hand-written tool loop",11,MUTE,sb=2),
    P("Haiku retrieval → Sonnet ranking",11,FG,font=MONO,sb=4)])
box(360,470,420,150, line=PURPLE, lw=2, anchor=MSO_ANCHOR.MIDDLE, paras=[
    P("Prospector",18,PURPLE_LT,b=True,font=HEAD),P("LangGraph",11,MUTE,sb=2),
    P("branching · human-in-the-loop",11,FG,sb=4)])

# Spine
box(800,180,430,690, fill="0E3144", line=BLUE, lw=2, radius=0.04, anchor=MSO_ANCHOR.TOP)
def stage(y,h,color,title,tag=None,emph=False):
    box(820,y,390,h, fill=(PANEL2 if emph else PANEL), line=color, lw=(3 if emph else 2),
        anchor=MSO_ANCHOR.MIDDLE, paras=[P(title,16,FG,b=True,font=HEAD)])
    if tag: label(1030,y+(h/2)-11,168,22,[P(tag,10,color if color!=GREEN else GREEN_LT,b=True,align=PP_ALIGN.RIGHT)])
stage(220,82,BLUE,"MCP tool belt",tag="retrieval")
stage(340,92,GREEN,"Deterministic scorer",tag="the oracle",emph=True)
stage(470,82,OCHRE,"Grounding contract")
stage(580,92,CLAY,"Guardrail veto",tag="rules decide")
label(820,710,390,26,[P("the model advises · the rules decide",13,BLUE_LT,b=True,align=PP_ALIGN.CENTER,font=HEAD)])
conn(1015,302,1015,340); conn(1015,432,1015,470); conn(1015,552,1015,580)

# Data sources (names only)
src=["USGS NWIS","USGS NLDI","NOAA","State ArcGIS","PAD-US","Postgres"]
gx=[1280,1570]; gy=[210,300,390]
for i,nm in enumerate(src):
    box(gx[i%2],gy[i//2],250,72, line=LINE, lw=1.4, anchor=MSO_ANCHOR.MIDDLE, paras=[P(nm,14,FG,b=True,font=HEAD,align=PP_ALIGN.CENTER)])
conn(1280,246,1212,265, color=BLUE, width=2.2)
label(1140,252,90,18,[P("fetch",10,BLUE_LT,align=PP_ALIGN.RIGHT)])

# Output
label(1560,478,300,22,[P("OUTPUT",11,FAINT,b=True,align=PP_ALIGN.CENTER,font=HEAD)])
box(1280,506,275,120, line=GREEN, lw=2, anchor=MSO_ANCHOR.MIDDLE, paras=[
    P("$0.02",34,GREEN_LT,b=True,align=PP_ALIGN.CENTER,font=HEAD),P("per decision",11,MUTE,align=PP_ALIGN.CENTER)])
box(1575,506,275,120, line=BLUE, lw=2, anchor=MSO_ANCHOR.MIDDLE, paras=[
    P("17–18 s",30,BLUE_LT,b=True,align=PP_ALIGN.CENTER,font=HEAD),P("latency",11,MUTE,align=PP_ALIGN.CENTER)])
label(1280,640,570,40,[P("→ ranked, grounded recommendations + guardrail verdicts, on the map",10.5,MUTE,align=PP_ALIGN.CENTER)])

# flow arrows
conn(300,270,358,270, color=BLUE, width=2.2)
conn(300,530,358,530, color=PURPLE, width=2.2)
conn(780,285,798,300, color=BLUE, width=2.2)
conn(780,525,798,470, color=PURPLE, width=2.2)
conn(1230,610,1278,566, color=CLAY, width=2.4)

# Eval band (one line)
band=box(70,900,1780,110, fill=PANEL_DK, line=BLUE, lw=1.6, radius=0.04, anchor=MSO_ANCHOR.TOP)
band.line._get_or_add_ln().append(band.line._get_or_add_ln().makeelement(qn("a:prstDash"),{"val":"dash"}))
label(92,924,1000,24,[P("OFFLINE EVAL HARNESS",11,BLUE_LT,b=True,font=HEAD)])
label(92,956,1700,28,[P("The scorer here is the eval oracle — same code · validated on 25 planner scenarios + an honest discovery eval",11,MUTE)])
conn(812,900,816,432, color=GREEN_LT, width=1.8, dashed=True)

path=os.path.join(OUT,"architecture_simple.pptx"); prs.save(path)
print("wrote", path, os.path.getsize(path), "bytes")
