#!/usr/bin/env python3
"""Batch 1 of the blank-space evidence cards, rebuilt as FULLY NATIVE pages so
every element stays editable in Canva (same pipeline as slides 5 & 8).

Each page reproduces the slide's CURRENT text verbatim (captured live) and adds
one real-data evidence card in the blank region:

  page 1 -> deck slide 10  "Real Bug Discovery"     + before/after code-diff card
  page 2 -> deck slide 13  "Measuring Discovery"     + discovery backtest scoreboard
  page 3 -> deck slide 12  "Prospector Mechanics"    + LangGraph pipeline diagram

Sources: agent/guardrails.py (_norm/ev_index fix), agent/eval/backtest_report.md,
agent/prospector_graph.py / prospector_nodes.py.
Output: deck/assets/fill_batch1.pptx
"""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn

OUT = os.path.join(os.path.dirname(__file__), "assets"); os.makedirs(OUT, exist_ok=True)
BG="2B2A45"; PANEL="38365C"; PANEL2="454270"; PANEL_DK="232239"; TERM="1E1D33"
YEL="F0C84C"; YEL_LT="F7DD85"; LAV="A89FE6"; LAV_LT="C5BEF2"
GREEN="5FB98C"; RED="E0685C"; RED_LT="EBA59C"; AMBER="E2A24A"
FG="FFFFFF"; CODE="E8E7F2"; MUTE="C2C0D6"; FAINT="8E8CAB"; LINE="55527E"
HEAD="Poppins"; BODY="Poppins"; MONO="Roboto Mono"
PXEMU=6350
def E(px): return Emu(int(round(px*PXEMU)))
def RGB(h): return RGBColor.from_string(h)

prs=Presentation(); prs.slide_width=E(1920); prs.slide_height=E(1080); BLANK=prs.slide_layouts[6]
def slide():
    s=prs.slides.add_slide(BLANK)
    bg=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,E(0),E(0),E(1920),E(1080))
    bg.fill.solid(); bg.fill.fore_color.rgb=RGB(BG); bg.line.fill.background(); bg.shadow.inherit=False
    return s
def P(t,size,color,b=False,font=BODY,align=PP_ALIGN.LEFT,sb=0,sa=1,ls=1.05):
    return {"t":t,"size":size,"color":color,"b":b,"font":font,"align":align,"sb":sb,"sa":sa,"ls":ls}
def MP(runs,align=PP_ALIGN.LEFT,sb=0,sa=1,ls=1.1):
    return {"runs":runs,"align":align,"sb":sb,"sa":sa,"ls":ls,"t":runs[0]["t"],"size":runs[0]["size"],"color":runs[0]["color"]}
def _tf(tf,paras,anchor,ml=14,mr=10,mt=6,mb=6):
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=E(ml); tf.margin_right=E(mr); tf.margin_top=E(mt); tf.margin_bottom=E(mb)
    for i,p in enumerate(paras):
        par=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        par.alignment=p.get("align",PP_ALIGN.LEFT); par.line_spacing=p.get("ls",1.05)
        par.space_after=Pt(p.get("sa",1)); par.space_before=Pt(p.get("sb",0))
        for run in p.get("runs",[p]):
            r=par.add_run(); r.text=run["t"]; f=r.font
            f.size=Pt(run["size"]); f.bold=run.get("b",False); f.name=run.get("font",BODY); f.color.rgb=RGB(run["color"])
def box(s,x,y,w,h,fill=PANEL,line=LINE,lw=1.6,radius=0.08,paras=None,anchor=MSO_ANCHOR.MIDDLE,ml=14,mr=10,mt=6,mb=6,dash=False):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,E(x),E(y),E(w),E(h))
    try: sp.adjustments[0]=radius
    except Exception: pass
    sp.fill.solid(); sp.fill.fore_color.rgb=RGB(fill)
    if line:
        sp.line.color.rgb=RGB(line); sp.line.width=Pt(lw)
        if dash: sp.line._get_or_add_ln().append(sp.line._get_or_add_ln().makeelement(qn("a:prstDash"),{"val":"dash"}))
    else: sp.line.fill.background()
    sp.shadow.inherit=False
    if paras: _tf(sp.text_frame,paras,anchor,ml=ml,mr=mr,mt=mt,mb=mb)
    return sp
def label(s,x,y,w,h,paras,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(E(x),E(y),E(w),E(h)); _tf(tb.text_frame,paras,anchor,ml=2,mr=2,mt=2,mb=2); return tb
def conn(s,x1,y1,x2,y2,color=LINE,width=1.4,dashed=False,arrow=True):
    c=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,E(x1),E(y1),E(x2),E(y2))
    c.line.color.rgb=RGB(color); c.line.width=Pt(width); ln=c.line._get_or_add_ln()
    if dashed: ln.append(ln.makeelement(qn("a:prstDash"),{"val":"dash"}))
    if arrow: ln.append(ln.makeelement(qn("a:tailEnd"),{"type":"triangle","w":"med","len":"med"}))
    return c

# ---- shared chrome ----
def footer(s): label(s,70,1016,400,28,[P("BlueLiner Agents",11,FAINT,font=HEAD)])
def text_chrome(s,title,subtitle,body_paras):
    label(s,70,66,800,230,[P(title,50,FG,b=True,font=HEAD,ls=1.0)])
    label(s,70,300,770,60,[P(subtitle,23,FG,font=HEAD)])
    label(s,976,70,874,900,body_paras); footer(s)
def bod(t): return P(t,16.5,FG,ls=1.28,sa=12)
def case_chrome(s,title,subtitle,sections):
    label(s,70,60,1340,120,[P(title,48,FG,b=True,font=HEAD,ls=1.0)])
    label(s,70,205,1340,56,[P(subtitle,22,FG,font=HEAD)])
    for (hy,htext,by,btext) in sections:
        label(s,70,hy,874,52,[P(htext,21,FG,b=True,font=HEAD)])
        label(s,70,by,874,120,[P(btext,14.5,MUTE,ls=1.3)])
    conn(s,1580,60,1580,1020,color=LINE,width=1.0,arrow=False)  # faint right divider
    footer(s)

# ================= CARD: before/after code diff (slide 10) =================
def bug_card(s,ox,oy,W=588,H=664):
    A=lambda x:ox+x; B=lambda y:oy+y
    box(s,A(0),B(0),W,H,fill=PANEL_DK,line=LINE,lw=1.5,radius=0.03)
    box(s,A(6),B(22),6,H-44,fill=YEL,line=None,radius=0.5)
    label(s,A(28),B(14),W-50,22,[P("THE BUG, AND THE ONE-LINE FIX",12,FAINT,b=True,font=HEAD)])
    label(s,A(28),B(40),W-50,22,[P("agent/guardrails.py",12,MUTE,font=MONO)])
    conn(s,A(24),B(72),A(W-24),B(72),color=LINE,arrow=False)
    # BEFORE
    label(s,A(26),B(86),W-50,22,[MP([{"t":"✗  ","size":13.5,"color":RED,"b":True},
        {"t":"BEFORE — guardrail fails open","size":13.5,"color":RED,"b":True}])])
    before=[MP([{"t":'ev = evidence.get(rec["river_id"])',"size":11.5,"color":CODE,"font":MONO}],sa=3),
            MP([{"t":'# model returns  "beaver_creek_md"',"size":11.5,"color":FAINT,"font":MONO}],sa=1),
            MP([{"t":'# evidence key is "beaver-creek-md"',"size":11.5,"color":FAINT,"font":MONO}],sa=1),
            MP([{"t":"# lookup ","size":11.5,"color":FAINT,"font":MONO},
                {"t":"MISSES","size":11.5,"color":RED,"font":MONO,"b":True},
                {"t":" -> veto never runs","size":11.5,"color":FAINT,"font":MONO}],sa=1)]
    box(s,A(24),B(116),W-48,128,fill=TERM,line=LINE,lw=1,radius=0.05,anchor=MSO_ANCHOR.TOP,ml=16,mt=12,paras=before)
    # AFTER
    label(s,A(26),B(262),W-50,22,[MP([{"t":"✓  ","size":13.5,"color":GREEN,"b":True},
        {"t":"AFTER — veto always fires","size":13.5,"color":GREEN,"b":True}])])
    after=[MP([{"t":"def _norm(rid):","size":11.5,"color":CODE,"font":MONO}],sa=3),
           MP([{"t":"    return (str(rid).strip().lower()","size":11.5,"color":CODE,"font":MONO}],sa=1),
           MP([{"t":'                   .replace("_","-"))',"size":11.5,"color":CODE,"font":MONO}],sa=1),
           MP([{"t":"ev_index = {_norm(k): (k, v) ...}","size":11.5,"color":CODE,"font":MONO}],sa=1),
           MP([{"t":"# canonical id ","size":11.5,"color":FAINT,"font":MONO},
               {"t":"-> guardrail matches","size":11.5,"color":GREEN,"font":MONO,"b":True}],sa=1)]
    box(s,A(24),B(292),W-48,150,fill=TERM,line=LINE,lw=1,radius=0.05,anchor=MSO_ANCHOR.TOP,ml=16,mt=12,paras=after)
    conn(s,A(24),B(470),A(W-24),B(470),color=LINE,arrow=False)
    box(s,A(24),B(488),W-48,72,fill=PANEL,line=GREEN,lw=1.4,radius=0.08,anchor=MSO_ANCHOR.MIDDLE,
        paras=[MP([{"t":"✓  ","size":13,"color":GREEN,"b":True},
                   {"t":"regression test asserts blocked water stays blocked","size":12.5,"color":FG}],align=PP_ALIGN.CENTER)])
    label(s,A(24),B(584),W-48,40,[P("a safety control that fails silently is worse than none",11.5,FAINT,align=PP_ALIGN.CENTER)])

def slide10(s):
    case_chrome(s,"Real Bug Discovery","A safety control that fails open is worse than none",
        [(352,"Safety Bypass",446,"The model reformatted a river ID, the safety lookup missed it, and the veto was silently skipped. The guardrail failed open."),
         (634,"Canonicalization Fix",728,"Canonicalize every ID before the veto runs, plus a regression test that asserts blocked water stays blocked.")])
    bug_card(s,966,312)

# ================= CARD: discovery backtest scoreboard (slide 13) =================
def discovery_card(s,ox,oy,W=812,H=560):
    A=lambda x:ox+x; B=lambda y:oy+y
    box(s,A(0),B(0),W,H,fill=PANEL,line=LINE,lw=1.5,radius=0.03)
    label(s,A(28),B(16),W-50,22,[P("DISCOVERY BACKTEST  ·  MD · VA · PA",11.5,FAINT,b=True,font=HEAD)])
    # two contrasting stat chips
    box(s,A(28),B(54),366,120,fill=PANEL_DK,line=GREEN,lw=1.6,radius=0.06,anchor=MSO_ANCHOR.TOP,mt=14,
        paras=[P("0.986",40,GREEN,b=True,align=PP_ALIGN.CENTER,font=HEAD),
               P("hard-neg AUC · topology",12,MUTE,align=PP_ALIGN.CENTER,sb=2),
               P("ranks hidden trout above near-trout decoys",10.5,FAINT,align=PP_ALIGN.CENTER,sb=2)])
    box(s,A(418),B(54),366,120,fill=PANEL_DK,line=AMBER,lw=1.6,radius=0.06,anchor=MSO_ANCHOR.TOP,mt=14,
        paras=[P("0.512",40,AMBER,b=True,align=PP_ALIGN.CENTER,font=HEAD),
               P("with public access enforced",12,MUTE,align=PP_ALIGN.CENTER,sb=2),
               P("~ chance: the access data gap",10.5,FAINT,align=PP_ALIGN.CENTER,sb=2)])
    label(s,A(28),B(188),W-50,20,[P("WHAT THE NUMBERS SAY",11,FAINT,b=True,font=HEAD)])
    label(s,A(28),B(212),W-56,30,[P("Topology alone separates held-out trout from near-trout decoys.",13,MUTE,ls=1.25)])
    label(s,A(28),B(238),W-56,52,[P("Demand verified public access and it falls to chance: the binding constraint is a data gap, not the model.",13,MUTE,ls=1.25)])
    box(s,A(28),B(300),W-56,58,fill=PANEL_DK,line=LINE,lw=1,radius=0.08,anchor=MSO_ANCHOR.MIDDLE,ml=18,
        paras=[MP([{"t":"Calibration: predicted confidence 0.6–0.8  →  ","size":12.5,"color":MUTE},
                   {"t":"53%","size":12.5,"color":GREEN,"b":True},
                   {"t":" actual held-out hit-rate","size":12.5,"color":MUTE}])])
    conn(s,A(28),B(386),A(W-28),B(386),color=LINE,arrow=False)
    label(s,A(28),B(400),W-50,22,[P("405 held-out reaches across 204 whole rivers, masked by flow path (levelpathid).",11.5,FAINT)])
    label(s,A(28),B(424),W-50,22,[P("Recall is a lower bound (positive-unlabeled); v1's 0.999 AUC was fake segment-adjacency.",11.5,FAINT)])

def slide13(s):
    text_chrome(s,"Measuring Discovery","An honest evaluation of results",[
        bod("An eval that refuses to flatter itself (MD, VA, and PA; 36,378 designated reaches)."),
        bod("Mask whole rivers by flow path, not random reaches. Otherwise discovery is just trivial segment in-painting (v1 scored a fake 0.999 AUC)."),
        bod("Hard-negative AUC: held-out trout versus only near-trout undesignated reaches. That's the honest, harder number."),
        bod("Positive-unlabeled: unlabeled is not the same as negative, so recall is a lower bound. The eval undercounts wins by construction."),
        bod("Finding: topology is a near-perfect lead generator, but the binding constraint, public access, is also the biggest data gap.")])
    discovery_card(s,70,402)

# ================= CARD: LangGraph pipeline diagram (slide 12) =================
def pipeline_card(s,ox,oy,W=812,H=576):
    A=lambda x:ox+x; B=lambda y:oy+y
    box(s,A(0),B(0),W,H,fill=PANEL,line=LINE,lw=1.5,radius=0.03)
    label(s,A(24),B(16),W-40,22,[P("PROSPECTOR  ·  LangGraph StateGraph",11.5,FAINT,b=True,font=HEAD)])
    nx,nw=46,330; y=52; step=58; nh=44
    nodes=[("generate_candidates",LINE,None),
           ("gather_evidence",LINE,None),
           ("infer_thermal",MUTE,"conditional branch · if ungauged"),
           ("score",LINE,None),
           ("reflect_verify",LINE,"drops private / low-confidence"),
           ("rank",YEL,"the only LLM call · writes rationale"),
           ("human_confirm",LAV,"interrupt + durable checkpoint"),
           ("update_flywheel",LINE,"records confirm/deny")]
    ys=[]
    for i,(name,col,tag) in enumerate(nodes):
        ny=y+i*step; ys.append(ny)
        dash=(name=="infer_thermal")
        box(s,A(nx),B(ny),nw,nh,fill=PANEL_DK,line=col,lw=(2.2 if col in (YEL,LAV) else 1.3),radius=0.12,
            anchor=MSO_ANCHOR.MIDDLE,ml=16,paras=[P(name,13.5,(YEL_LT if col==YEL else LAV_LT if col==LAV else FG),
                                                     b=(col in (YEL,LAV)),font=MONO)])
        if tag:
            tc=YEL_LT if col==YEL else (LAV_LT if col==LAV else FAINT)
            label(s,A(nx+nw+18),B(ny+11),W-(nx+nw+40),22,[P(tag,11,tc)])
    for i in range(len(nodes)-1):
        conn(s,A(nx+nw/2),B(ys[i]+nh),A(nx+nw/2),B(ys[i+1]),color=LINE,width=1.3,arrow=True)
    label(s,A(24),B(H-40),W-40,30,[P("deterministic spine; the LLM touches only the top-K shortlist.",11.5,FAINT,align=PP_ALIGN.CENTER)])

def slide12(s):
    text_chrome(s,"Prospector Mechanics","How the Prospector ranks",[
        bod("Cheap deterministic ranking at scale; the LLM only where it adds value."),
        bod("100K+ reaches ranked by topology proximity to known trout water (a geometry proxy), plus flow, thermal, and access signals."),
        bod("Deterministic confidence; the LLM writes the rationale on the shortlist only."),
        bod("A human-confirm interrupt with durable checkpointing is wired to feed a calibration flywheel as anglers confirm picks; it is proven on a captured resume trace and runs headless in the live demo. Deterministic ranks 60K reaches for free; the LLM touches only the top-K.")])
    pipeline_card(s,70,398)

slide10(slide()); slide13(slide()); slide12(slide())
path=os.path.join(OUT,"fill_batch1.pptx"); prs.save(path)
print("wrote",path,os.path.getsize(path),"bytes,",len(prs.slides._sldIdLst),"slides")
