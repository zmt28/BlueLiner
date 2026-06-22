#!/usr/bin/env python3
"""Slides 5 & 8 rebuilt as FULLY NATIVE, EDITABLE PowerPoint pages so that
EVERYTHING (title, subtitle, body, footer, AND the evidence cards) becomes an
editable Canva element after import -- no flattened images.

 - Slide 5  "Operating Principle"  -> trace card (one decision, fully traced)
 - Slide 8  "How I Measured It"     -> eval scoreboard (25 scenarios)

Card content is the SAME real agent/eval output as deck/build_cards.py, but
drawn with native shapes/text in the deck palette (indigo + yellow + green/red).
Output: deck/assets/slides_5_8.pptx  (2 slides), for import-design-from-url.
"""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn

OUT = os.path.join(os.path.dirname(__file__), "assets"); os.makedirs(OUT, exist_ok=True)
# ---- deck-matched palette (identical to build_deck_graphics.py) ----
BG="2B2A45"; PANEL="38365C"; PANEL2="454270"; PANEL_DK="232239"; TERM="1E1D33"
YEL="F0C84C"; YEL_LT="F7DD85"; LAV="A89FE6"; LAV_LT="C5BEF2"
GREEN="5FB98C"; RED="E0685C"; RED_LT="EBA59C"
FG="FFFFFF"; MUTE="C2C0D6"; FAINT="8E8CAB"; LINE="55527E"
HEAD="Poppins"; BODY="Poppins"; MONO="Roboto Mono"
PXEMU=6350
def E(px): return Emu(int(round(px*PXEMU)))
def RGB(h): return RGBColor.from_string(h)

prs=Presentation(); prs.slide_width=E(1920); prs.slide_height=E(1080); BLANK=prs.slide_layouts[6]
def slide():
    s=prs.slides.add_slide(BLANK)
    bg=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,E(0),E(0),E(1920),E(1080))
    bg.fill.solid(); bg.fill.fore_color.rgb=RGB(BG); bg.line.fill.background(); bg.shadow.inherit=False
    return s
def P(t,size,color,b=False,font=BODY,align=PP_ALIGN.LEFT,sb=0,sa=1,ls=1.05):
    return {"t":t,"size":size,"color":color,"b":b,"font":font,"align":align,"sb":sb,"sa":sa,"ls":ls}
def _tf(tf,paras,anchor,ml=14,mr=10,mt=6,mb=6):
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=E(ml); tf.margin_right=E(mr); tf.margin_top=E(mt); tf.margin_bottom=E(mb)
    for i,p in enumerate(paras):
        par=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        par.alignment=p.get("align",PP_ALIGN.LEFT); par.line_spacing=p.get("ls",1.05)
        par.space_after=Pt(p.get("sa",1)); par.space_before=Pt(p.get("sb",0))
        for run in p.get("runs",[p]):
            r=par.add_run(); r.text=run["t"]; f=r.font
            f.size=Pt(run["size"]); f.bold=run.get("b",False); f.name=run.get("font",BODY); f.color.rgb=RGB(run["color"])
def MP(runs,align=PP_ALIGN.LEFT,sb=0,sa=1,ls=1.05):
    """multi-run paragraph"""
    return {"runs":runs,"align":align,"sb":sb,"sa":sa,"ls":ls,"t":runs[0]["t"],"size":runs[0]["size"],"color":runs[0]["color"]}
def box(s,x,y,w,h,fill=PANEL,line=LINE,lw=1.6,radius=0.08,paras=None,anchor=MSO_ANCHOR.MIDDLE,ml=14,mr=10,mt=6,mb=6,shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp=s.shapes.add_shape(shape,E(x),E(y),E(w),E(h))
    try: sp.adjustments[0]=radius
    except Exception: pass
    sp.fill.solid(); sp.fill.fore_color.rgb=RGB(fill)
    if line: sp.line.color.rgb=RGB(line); sp.line.width=Pt(lw)
    else: sp.line.fill.background()
    sp.shadow.inherit=False
    if paras: _tf(sp.text_frame,paras,anchor,ml=ml,mr=mr,mt=mt,mb=mb)
    return sp
def label(s,x,y,w,h,paras,anchor=MSO_ANCHOR.TOP,align=None):
    tb=s.shapes.add_textbox(E(x),E(y),E(w),E(h)); _tf(tb.text_frame,paras,anchor,ml=2,mr=2,mt=2,mb=2); return tb
def conn(s,x1,y1,x2,y2,color=LINE,width=1.2,dashed=False,arrow=False):
    c=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,E(x1),E(y1),E(x2),E(y2))
    c.line.color.rgb=RGB(color); c.line.width=Pt(width); ln=c.line._get_or_add_ln()
    if dashed: ln.append(ln.makeelement(qn("a:prstDash"),{"val":"dash"}))
    if arrow: ln.append(ln.makeelement(qn("a:tailEnd"),{"type":"triangle","w":"med","len":"med"}))
    return c

# ---- shared page chrome: title / subtitle / body / footer ----
def chrome(s,title,subtitle,body_paras):
    label(s,70,66,800,230,[P(title,50,FG,b=True,font=HEAD,ls=1.0)])
    label(s,70,300,770,60,[P(subtitle,23,FG,font=HEAD)])
    label(s,976,70,874,900,body_paras)
    label(s,70,1016,400,28,[P("BlueLiner Agents",11,FAINT,font=HEAD)])
def bod(text): return P(text,16.5,FG,ls=1.28,sa=12)

# ======================= SLIDE 5: trace card =======================
def trace_card(s,ox,oy):
    A=lambda lx:ox+lx; B=lambda ly:oy+ly
    box(s,A(0),B(0),812,594,fill=PANEL_DK,line=LINE,lw=1.5,radius=0.03)          # panel
    box(s,A(6),B(22),6,486,fill=YEL,line=None,radius=0.5)                        # accent rail
    label(s,A(28),B(14),620,22,[P("ONE DECISION, FULLY TRACED",12,FAINT,b=True,font=HEAD)])
    label(s,A(28),B(38),756,28,[P("“Beaver Creek near Hagerstown: worth it this weekend?”",14,FG,font=HEAD)])
    label(s,A(28),B(70),756,22,[P("Haiku drives retrieval  ·  Sonnet writes the ranking",12.5,MUTE)])
    conn(s,A(28),B(104),A(784),B(104),color=LINE)
    label(s,A(28),B(112),420,20,[P("TOOLS CALLED, IN ORDER",11,FAINT,b=True,font=HEAD)])
    tools=[("1","get_candidate_rivers",""),("2","get_forecast","NOAA · live"),
           ("3","get_river_conditions","2 rivers"),("4","get_access","2 rivers"),
           ("5","get_user_catch_history","")]
    ly=142
    for n,name,desc in tools:
        label(s,A(34),B(ly),330,24,[MP([{"t":n+"   ","size":13.5,"color":YEL,"b":True,"font":MONO},
                                        {"t":name,"size":13.5,"color":FG,"font":MONO}])])
        if desc: label(s,A(360),B(ly),330,24,[P(desc,12,MUTE,font=MONO)])
        ly+=28
    conn(s,A(28),B(296),A(784),B(296),color=LINE)
    label(s,A(28),B(310),32,26,[P("✗",16,RED,b=True)])
    label(s,A(58),B(312),724,24,[P("beaver-creek-md   blocked: private-only access, no public entry",13,RED_LT)])
    label(s,A(28),B(350),32,28,[P("✓",16,GREEN,b=True)])
    label(s,A(58),B(350),724,28,[P("Gunpowder Falls   (green)",15,GREEN,b=True,font=HEAD)])
    label(s,A(58),B(382),724,22,[P("56°F  ·  flow 1.09x median  ·  6 public access points",12.5,MUTE)])
    label(s,A(58),B(406),724,22,[P("within your proven brown + rainbow trout bands; fresh reading",12,FAINT)])
    conn(s,A(16),B(540),A(796),B(540),color=LINE)
    label(s,A(0),B(554),812,26,[MP([{"t":"grounding ok ","size":12.5,"color":MUTE,"font":MONO},
                                    {"t":"✓","size":12.5,"color":GREEN,"font":MONO,"b":True},
                                    {"t":"    unsourced: none     ·     16.6 s     ·     $0.023","size":12.5,"color":MUTE,"font":MONO}],
                                   align=PP_ALIGN.CENTER)])

def slide5():
    s=slide()
    chrome(s,"Operating Principle","Clarity over complexity",[
        bod("Legibility over cleverness. It's optimized for explainability end to end."),
        bod("Inspectable: it's a hand-written tool loop, not a black box, so you can walk through every step."),
        bod("Single source of truth: the deterministic scorer is both the agent's tool and the eval's oracle, the same code."),
        bod("Visible evolution: the v0 to v3 progression lives in git, and I show what didn't work.")])
    trace_card(s,70,400)

# ======================= SLIDE 8: eval scoreboard =======================
def eval_card(s,ox,oy):
    A=lambda lx:ox+lx; B=lambda ly:oy+ly
    box(s,A(0),B(0),812,594,fill=PANEL,line=LINE,lw=1.5,radius=0.03)
    label(s,A(28),B(14),520,22,[P("v3 EVAL  ·  25 SCENARIOS",12,FAINT,b=True,font=HEAD)])
    cw=366; gap=24; x1=28; x2=28+cw+gap
    for cx,lab in [(x1,"safety violations"),(x2,"hallucinated readings")]:
        box(s,A(cx),B(48),cw,92,fill=PANEL_DK,line=GREEN,lw=1.6,radius=0.07,anchor=MSO_ANCHOR.MIDDLE,paras=[
            P("0",42,GREEN,b=True,align=PP_ALIGN.CENTER,font=HEAD),
            P(lab,12.5,MUTE,align=PP_ALIGN.CENTER,sb=2)])
    label(s,A(28),B(156),600,20,[P("TOP-1 AGREEMENT, BY SCENARIO TYPE",11,FAINT,b=True,font=HEAD)])
    cats=[("adversarial","memory"),("all-blocked","private"),("flood","stale"),
          ("ideal","tie"),("marginal","too-warm")]
    y=186
    for leftc,rightc in cats:
        for cx,name in [(x1,leftc),(x2,rightc)]:
            box(s,A(cx),B(y),cw,48,fill=PANEL_DK,line=LINE,lw=1,radius=0.12,anchor=MSO_ANCHOR.MIDDLE,
                ml=20,paras=[P(name,13.5,FG)])
            label(s,A(cx+cw-130),B(y+13),118,24,[MP([{"t":"100% ","size":13.5,"color":GREEN,"b":True},
                                                     {"t":"✓","size":13.5,"color":GREEN,"b":True}],align=PP_ALIGN.RIGHT)])
        y+=58
    label(s,A(28),B(508),760,44,[
        P("Oracle: BlueLiner's own deterministic scorer, parity-tested.",11.5,FAINT,sa=4),
        P("Agreement = top pick is safe and in the oracle's best-rated tier.",11.5,FAINT)])

def slide8():
    s=slide()
    chrome(s,"How I Measured It","The Eval Is the Product",[
        bod("Evaluated against a deterministic oracle: BlueLiner's own scorer, parity-tested in 840 cases against production."),
        bod("25 scenarios: ideal, flood, too-warm, private, stale, adversarial, memory, ties, and all-blocked."),
        bod("Honest caveat: safety at 0% in v1 and v2 is luck. Only v3 cannot recommend blocked water by construction."),
        bod("Honest caveat: personalization is confounded (n=4), so the durable signal is qualitative.")])
    eval_card(s,70,400)

slide5(); slide8()
path=os.path.join(OUT,"slides_5_8.pptx"); prs.save(path)
print("wrote",path,os.path.getsize(path),"bytes,",len(prs.slides._sldIdLst),"slides")
