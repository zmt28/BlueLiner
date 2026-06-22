#!/usr/bin/env python3
"""Build the three BlueLiner deck graphics as NATIVE, EDITABLE PowerPoint shapes.

One shape per box with its text inside it, real connectors/arrows, standalone
text labels — so after importing into Canva every element is editable (drag,
retype, recolor) with no code. 16:9 slides at 1920x1080 proportions.

Output: deck/assets/blueliner_graphics.pptx  (3 slides)
"""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn

OUT = os.path.join(os.path.dirname(__file__), "assets")
os.makedirs(OUT, exist_ok=True)

# palette
NAVY="0B2A3A"; PANEL="102E41"; PANEL2="143A50"; PANEL_DK="0C2535"
BLUE="5BA8C8"; BLUE_LT="95C5D9"; GREEN="4A8C5C"; GREEN_LT="7FBE8E"
OCHRE="B7892F"; CLAY="B3473B"; CLAY_LT="E0A59B"; PURPLE="7A3DB8"; PURPLE_LT="A878DA"
FG="EAF2F6"; MUTE="9DB6C2"; FAINT="6E8C9C"; LINE="3E5C6E"
HEAD="Poppins"; BODY="Poppins"; MONO="Roboto Mono"

PXEMU = 6350  # EMU per px on a 1920px (=13.333in) wide slide
def E(px): return Emu(int(round(px*PXEMU)))
def RGB(h): return RGBColor.from_string(h)

prs = Presentation()
prs.slide_width = E(1920); prs.slide_height = E(1080)
BLANK = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, E(0),E(0),E(1920),E(1080))
    bg.fill.solid(); bg.fill.fore_color.rgb=RGB(NAVY); bg.line.fill.background(); bg.shadow.inherit=False
    return s

def _fill_tf(tf, paras, anchor, ml=14, mr=10, mt=8, mb=6):
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=E(ml); tf.margin_right=E(mr); tf.margin_top=E(mt); tf.margin_bottom=E(mb)
    for i,p in enumerate(paras):
        para = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        para.alignment = p.get("align", PP_ALIGN.LEFT)
        para.line_spacing = p.get("ls", 1.05)
        para.space_after = Pt(p.get("sa", 1))
        para.space_before = Pt(p.get("sb", 0))
        run = para.add_run(); run.text = p["t"]
        f = run.font
        f.size=Pt(p["size"]); f.bold=p.get("b",False); f.name=p.get("font",BODY)
        f.color.rgb=RGB(p["color"])

def box(s, x,y,w,h, fill=PANEL, line=LINE, lw=1.6, radius=0.08, paras=None,
        anchor=MSO_ANCHOR.TOP):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, E(x),E(y),E(w),E(h))
    try: sp.adjustments[0]=radius
    except Exception: pass
    sp.fill.solid(); sp.fill.fore_color.rgb=RGB(fill)
    if line: sp.line.color.rgb=RGB(line); sp.line.width=Pt(lw)
    else: sp.line.fill.background()
    sp.shadow.inherit=False
    if paras: _fill_tf(sp.text_frame, paras, anchor)
    return sp

def label(s, x,y,w,h, paras, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(E(x),E(y),E(w),E(h))
    _fill_tf(tb.text_frame, paras, anchor, ml=2, mr=2, mt=2, mb=2)
    return tb

def conn(s, x1,y1,x2,y2, color=MUTE, width=2.0, dashed=False, arrow=True):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, E(x1),E(y1),E(x2),E(y2))
    c.line.color.rgb=RGB(color); c.line.width=Pt(width)
    ln = c.line._get_or_add_ln()
    if dashed: ln.append(ln.makeelement(qn("a:prstDash"), {"val":"dash"}))
    if arrow:  ln.append(ln.makeelement(qn("a:tailEnd"), {"type":"triangle","w":"med","len":"med"}))
    return c

def hdr(s, cx, text):
    label(s, cx-200, 132, 400, 30, [{"t":text,"size":11,"b":True,"color":FAINT,"align":PP_ALIGN.CENTER,"font":HEAD}])

# small helpers for paragraph dicts
def P(t,size,color,b=False,font=BODY,align=PP_ALIGN.LEFT,sb=0,sa=1,ls=1.05):
    return {"t":t,"size":size,"color":color,"b":b,"font":font,"align":align,"sb":sb,"sa":sa,"ls":ls}

# =====================================================================
# SLIDE 1 — Architecture
# =====================================================================
def architecture():
    s = slide()
    label(s, 70, 36, 1780, 70, [P("Two Agents, One Trustworthy Spine", 32, FG, b=True, font=HEAD)])
    hdr(s,175,"REQUESTS"); hdr(s,520,"ORCHESTRATION"); hdr(s,1015,"SHARED TRUSTWORTHY SPINE"); hdr(s,1480,"DATA SOURCES")

    # Requests
    box(s,70,180,230,150, line=BLUE, lw=2, paras=[
        P("Trip Planner",13,BLUE_LT,b=True,font=HEAD),
        P("“Where should I fish this weekend — and is it worth the drive?”",10,FG,sb=2)])
    box(s,70,470,230,150, line=PURPLE, lw=2, paras=[
        P("Prospector",13,PURPLE_LT,b=True,font=HEAD),
        P("“Find undesignated but fishable trout water”",10,FG,sb=2)])

    # Orchestration
    box(s,340,180,420,220, line=BLUE, lw=2, paras=[
        P("Trip Planner",16,BLUE_LT,b=True,font=HEAD),
        P("hand-written tool loop",10,MUTE,sa=4),
        P("Haiku → cheap, tool-heavy retrieval loop",10.5,FG,font=MONO),
        P("Sonnet → final ranking",10.5,FG,font=MONO,sa=4)])
    box(s,560,348,168,30, fill=PANEL_DK, line=BLUE, lw=1, radius=0.5, paras=[
        P("17 lines orchestration",9,BLUE_LT,align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    box(s,340,450,420,250, line=PURPLE, lw=2, paras=[
        P("Prospector",16,PURPLE_LT,b=True,font=HEAD),
        P("LangGraph",10,MUTE,sa=4),
        P("branching state machine",10.5,FG),
        P("human-in-the-loop confirm",10.5,FG),
        P("durable checkpoints · interrupt",10.5,FG,sa=4)])
    box(s,560,648,150,30, fill=PANEL_DK, line=PURPLE, lw=1, radius=0.5, paras=[
        P("38 lines (2.2×)",9,PURPLE_LT,align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)

    # Spine container + stages
    box(s,800,168,430,724, fill="0E3144", line=BLUE, lw=2, radius=0.04)
    box(s,820,198,390,92, line=BLUE, lw=2, paras=[
        P("MCP tool belt",13,FG,b=True,font=HEAD),
        P("conditions · gauges · 30-yr medians",9,MUTE,sb=1),
        P("trout designations · access · catch-log memory",9,MUTE)])
    label(s,1090,206,118,20,[P("retrieval",9,BLUE_LT,b=True,align=PP_ALIGN.RIGHT)])
    box(s,820,312,390,112, fill=PANEL2, line=GREEN, lw=3, paras=[
        P("Deterministic scorer",13,FG,b=True,font=HEAD),
        P("single source of truth",9,MUTE,sb=1),
        P("water-temp band + flow-vs-median",9,MUTE),
        P("parity-tested 840 cases against production",8.5,GREEN_LT,sb=2)])
    label(s,1090,320,118,20,[P("the oracle",9,GREEN_LT,b=True,align=PP_ALIGN.RIGHT)])
    box(s,820,446,390,92, line=OCHRE, lw=2, paras=[
        P("Grounding contract",13,FG,b=True,font=HEAD),
        P("every number must trace to a tool result —",9,MUTE,sb=1),
        P("else regenerate once, then strip it",9,MUTE)])
    box(s,820,560,390,124, line=CLAY, lw=2, paras=[
        P("Guardrail veto",13,FG,b=True,font=HEAD),
        P("flood (flow >3× median) · trout-ethics temp band",9,MUTE,sb=1),
        P("private-access block · staleness demotion",9,MUTE),
        P("IDs canonicalized before the veto — the bug fix",8.5,CLAY_LT,sb=2)])
    label(s,1090,568,118,20,[P("rules decide",9,CLAY_LT,b=True,align=PP_ALIGN.RIGHT)])
    label(s,820,700,390,24,[P("the model advises · the rules decide",11,BLUE_LT,b=True,align=PP_ALIGN.CENTER,font=HEAD)])
    # down arrows in spine
    conn(s,1015,290,1015,312); conn(s,1015,424,1015,446); conn(s,1015,538,1015,560)

    # Data sources grid
    src=[("USGS NWIS","flow + temp · IV + daily"),("USGS NLDI","topology · COMID"),
         ("NOAA","weather enrichment"),("State ArcGIS","trout designations"),
         ("PAD-US","public access / lands"),("Postgres","catch-log memory")]
    gx=[1270,1560]; gy=[180,272,364]
    for i,(nm,sub) in enumerate(src):
        box(s,gx[i%2],gy[i//2],250,80, line=LINE, lw=1.4, paras=[
            P(nm,12,FG,b=True,font=HEAD), P(sub,9.5,MUTE,sb=1)])
    conn(s,1270,225,1212,250, color=BLUE, width=2.2)
    label(s,1120,232,90,18,[P("fetch",9,BLUE_LT,align=PP_ALIGN.RIGHT)])

    # Output
    label(s,1480,452,300,22,[P("OUTPUT",11,FAINT,b=True,align=PP_ALIGN.CENTER,font=HEAD)])
    box(s,1270,478,260,210, line=LINE, lw=1.4, anchor=MSO_ANCHOR.MIDDLE, paras=[
        P("Delivered",14,FG,b=True,align=PP_ALIGN.CENTER,font=HEAD,sa=4),
        P("ranked recommendations",10,MUTE,align=PP_ALIGN.CENTER),
        P("+ grounded citations",10,MUTE,align=PP_ALIGN.CENTER),
        P("+ guardrail verdicts",10,MUTE,align=PP_ALIGN.CENTER),
        P("rendered on the map",9.5,FAINT,align=PP_ALIGN.CENTER,sb=6)])
    box(s,1560,478,290,100, line=GREEN, lw=2, anchor=MSO_ANCHOR.MIDDLE, paras=[
        P("$0.02",30,GREEN_LT,b=True,align=PP_ALIGN.CENTER,font=HEAD),
        P("per decision",9.5,MUTE,align=PP_ALIGN.CENTER)])
    box(s,1560,588,290,100, line=BLUE, lw=2, anchor=MSO_ANCHOR.MIDDLE, paras=[
        P("17–18 s",27,BLUE_LT,b=True,align=PP_ALIGN.CENTER,font=HEAD),
        P("end-to-end latency",9.5,MUTE,align=PP_ALIGN.CENTER)])

    # flow arrows
    conn(s,300,250,338,250, color=BLUE, width=2.2)
    conn(s,300,540,338,540, color=PURPLE, width=2.2)
    conn(s,760,260,798,290, color=BLUE, width=2.2)
    conn(s,760,540,798,440, color=PURPLE, width=2.2)
    conn(s,1230,620,1268,580, color=CLAY, width=2.4)

    # Eval band
    band = box(s,70,904,1780,118, fill=PANEL_DK, line=BLUE, lw=1.6, radius=0.04)
    band.line._get_or_add_ln().append(band.line._get_or_add_ln().makeelement(qn("a:prstDash"),{"val":"dash"}))
    label(s,92,916,900,24,[P("OFFLINE EVAL HARNESS",11,BLUE_LT,b=True,font=HEAD)])
    label(s,92,946,1700,50,[
        P("Planner — 25 scenarios: ideal · flood · too-warm · private · stale · adversarial · memory · ties · all-blocked",9.5,MUTE),
        P("Discovery — flow-path masking · hard-negative AUC · positive-unlabeled (recall = lower bound)",9.5,MUTE,sb=2)])
    label(s,930,916,640,24,[P("↑ the scorer here is the eval oracle — same code",9.5,GREEN_LT)])
    conn(s,812,904,816,424, color=GREEN_LT, width=1.8, dashed=True)

    # legend
    leg=[(BLUE,"Trip Planner"),(PURPLE,"Prospector"),(GREEN,"trusted / oracle"),(OCHRE,"grounding"),(CLAY,"guardrail / block")]
    lx=560
    for color,text in leg:
        sw=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,E(lx),E(1040),E(16),E(16))
        sw.adjustments[0]=0.3; sw.fill.solid(); sw.fill.fore_color.rgb=RGB(color); sw.line.fill.background(); sw.shadow.inherit=False
        label(s,lx+22,1037,180,22,[P(text,9.5,MUTE)])
        lx += 40 + len(text)*7.0

# =====================================================================
# SLIDE 2 — v0->v3 staircase
# =====================================================================
def staircase():
    s = slide()
    label(s,70,36,1780,50,[P("Trip Planner: The v0 → v3 Staircase",32,FG,b=True,font=HEAD)])
    label(s,70,96,1780,40,[P("Grounding + guardrails turn a confident liar into a trustworthy assistant (25 scenarios)",15,MUTE)])
    cols=[("v0","naive prompt, no tools",CLAY),("v1","tool-grounded",BLUE),
          ("v2","+ catch-log memory",OCHRE),("v3","+ guardrails & grounding",GREEN)]
    rows=[("Recommendation agreement","higher is better",["8%","100%","100%","100%"],"high"),
          ("Safety violations","lower is better",["16%","0%","0%","0%"],"low"),
          ("Hallucinated readings","lower is better",["100%","4%","12%","0%"],"low")]
    LBLW=330; x0=70; gap=18
    cw=(1850-x0-LBLW-gap*4)/4
    top=170; headh=104; rowh=150; rgap=14
    for i,(v,sub,color) in enumerate(cols):
        cx=x0+LBLW+i*(cw+gap); emph=(v=="v3")
        box(s,cx,top,cw,headh, fill=(PANEL2 if emph else PANEL), line=color, lw=(3 if emph else 2),
            anchor=MSO_ANCHOR.MIDDLE, paras=[
            P(v,30,color,b=True,align=PP_ALIGN.CENTER,font=HEAD),
            P(sub,11,(FG if emph else MUTE),align=PP_ALIGN.CENTER)])
    ry=top+headh+22
    for ri,(lab,note,vals,good) in enumerate(rows):
        yy=ry+ri*(rowh+rgap)
        label(s,x0,yy,LBLW-10,rowh,[P(lab,17,FG,b=True,font=HEAD),P(note,11,FAINT,sb=2)],anchor=MSO_ANCHOR.MIDDLE)
        for ci,(v,sub,color) in enumerate(cols):
            cx=x0+LBLW+ci*(cw+gap); val=vals[ci]; num=float(val.replace("%",""))
            if good=="high": cell=GREEN if num>=90 else (OCHRE if num>=50 else CLAY)
            else: cell=GREEN if num==0 else (OCHRE if num<=15 else CLAY)
            box(s,cx,yy,cw,rowh, line=cell, lw=(2.6 if v=="v3" else 1.6), anchor=MSO_ANCHOR.MIDDLE,
                paras=[P(val,40,cell,b=True,align=PP_ALIGN.CENTER,font=HEAD)])
    label(s,70,1004,1780,30,[P("v2’s hallucination bump (4% → 12%) is real — memory added unsourced numbers; the v3 grounding contract drove it to 0%.",14,MUTE)])

# =====================================================================
# SLIDE 3 — Orchestration A/B
# =====================================================================
def orchestration_ab():
    s = slide()
    label(s,70,36,1780,50,[P("Engineering Judgment: Right Tool for the Job",32,FG,b=True,font=HEAD)])
    label(s,70,150,1780,30,[P("SAME v3 PLANNER · 25 SCENARIOS · ONLY ORCHESTRATION CHANGES",13,FAINT,b=True,align=PP_ALIGN.CENTER,font=HEAD)])
    cards=[(360,"Hand-written loop",BLUE,BLUE_LT,"17","lines of orchestration","linear planner"),
           (1010,"LangGraph",PURPLE,PURPLE_LT,"38","lines (2.2×)","branching + HITL")]
    cy=200; ch=420; cw=560
    for x,name,color,lt,lines,lines_sub,foot in cards:
        box(s,x,cy,cw,ch, line=color, lw=2.4, radius=0.05)
        label(s,x,cy+28,cw,50,[P(name,24,lt,b=True,align=PP_ALIGN.CENTER,font=HEAD)])
        label(s,x,cy+92,cw,90,[P("100%",60,GREEN_LT,b=True,align=PP_ALIGN.CENTER,font=HEAD)])
        label(s,x,cy+188,cw,28,[P("scenario quality",13,MUTE,align=PP_ALIGN.CENTER)])
        conn(s,x+50,cy+232,x+cw-50,cy+232, color=LINE, width=1.4, arrow=False)
        label(s,x,cy+248,cw,80,[P(lines,52,lt,b=True,align=PP_ALIGN.CENTER,font=MONO)])
        label(s,x,cy+340,cw,28,[P(lines_sub,13,MUTE,align=PP_ALIGN.CENTER)])
        box(s,x+cw/2-95,cy+372,190,32, fill=PANEL_DK, line=color, lw=1, radius=0.5,
            anchor=MSO_ANCHOR.MIDDLE, paras=[P(foot,12,lt,align=PP_ALIGN.CENTER)])
    label(s,920,290,80,60,[P("=",40,MUTE,b=True,align=PP_ALIGN.CENTER,font=HEAD)])
    label(s,920,452,80,40,[P("vs",22,MUTE,b=True,align=PP_ALIGN.CENTER,font=HEAD)])
    box(s,360,664,1210,360-20, line=LINE, lw=1.5, fill=PANEL_DK, radius=0.04, paras=[
        P("Decision",16,FG,b=True,font=HEAD,sa=6),
        P("Hand-loop for the linear Trip Planner — less code, fully legible.",13.5,MUTE),
        P("LangGraph for the branching, human-in-the-loop Prospector — where interrupt + durable checkpoints actually earn their cost.",13.5,MUTE,sa=8),
        P("Frameworks are not a quality lever — claiming so would be a confound.",13.5,BLUE_LT,b=True)])

architecture(); staircase(); orchestration_ab()
path=os.path.join(OUT,"blueliner_graphics.pptx"); prs.save(path)
print("wrote", path, os.path.getsize(path), "bytes,", len(prs.slides._sldIdLst), "slides")
