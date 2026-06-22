#!/usr/bin/env python3
"""Deck-matched, NATIVE editable graphics (3 slides) for merge into the deck.
Palette: indigo bg + yellow primary accent + white text; lavender as the
secondary (Prospector) hue; green/red kept ONLY for good/bad data on the
staircase. One shape per node with text inside -> editable Canva elements.
Output: deck/assets/deck_graphics.pptx
"""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn

OUT = os.path.join(os.path.dirname(__file__), "assets"); os.makedirs(OUT, exist_ok=True)
# ---- deck-matched palette ----
BG="2B2A45"; PANEL="38365C"; PANEL2="454270"; PANEL_DK="232239"
YEL="F0C84C"; YEL_LT="F7DD85"; LAV="A89FE6"; LAV_LT="C5BEF2"
GREEN="5FB98C"; RED="E0685C"
FG="FFFFFF"; MUTE="C2C0D6"; FAINT="8E8CAB"; LINE="55527E"
HEAD="Poppins"; BODY="Poppins"; MONO="Roboto Mono"
PXEMU=6350
def E(px): return Emu(int(round(px*PXEMU)))
def RGB(h): return RGBColor.from_string(h)

prs=Presentation(); prs.slide_width=E(1920); prs.slide_height=E(1080); BLANK=prs.slide_layouts[6]
def slide():
    s=prs.slides.add_slide(BLANK)
    bg=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,E(0),E(0),E(1920),E(1080))
    bg.fill.solid(); bg.fill.fore_color.rgb=RGB(BG); bg.line.fill.background(); bg.shadow.inherit=False
    return s
def P(t,size,color,b=False,font=BODY,align=PP_ALIGN.LEFT,sb=0,sa=1,ls=1.05):
    return {"t":t,"size":size,"color":color,"b":b,"font":font,"align":align,"sb":sb,"sa":sa,"ls":ls}
def _tf(tf,paras,anchor,ml=14,mr=10,mt=8,mb=6):
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=E(ml); tf.margin_right=E(mr); tf.margin_top=E(mt); tf.margin_bottom=E(mb)
    for i,p in enumerate(paras):
        par=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        par.alignment=p.get("align",PP_ALIGN.LEFT); par.line_spacing=p.get("ls",1.05)
        par.space_after=Pt(p.get("sa",1)); par.space_before=Pt(p.get("sb",0))
        r=par.add_run(); r.text=p["t"]; f=r.font
        f.size=Pt(p["size"]); f.bold=p.get("b",False); f.name=p.get("font",BODY); f.color.rgb=RGB(p["color"])
def box(s,x,y,w,h,fill=PANEL,line=LINE,lw=1.6,radius=0.08,paras=None,anchor=MSO_ANCHOR.MIDDLE):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,E(x),E(y),E(w),E(h))
    try: sp.adjustments[0]=radius
    except Exception: pass
    sp.fill.solid(); sp.fill.fore_color.rgb=RGB(fill)
    if line: sp.line.color.rgb=RGB(line); sp.line.width=Pt(lw)
    else: sp.line.fill.background()
    sp.shadow.inherit=False
    if paras: _tf(sp.text_frame,paras,anchor)
    return sp
def label(s,x,y,w,h,paras,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(E(x),E(y),E(w),E(h)); _tf(tb.text_frame,paras,anchor,ml=2,mr=2,mt=2,mb=2); return tb
def conn(s,x1,y1,x2,y2,color=MUTE,width=2.0,dashed=False,arrow=True):
    c=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,E(x1),E(y1),E(x2),E(y2))
    c.line.color.rgb=RGB(color); c.line.width=Pt(width); ln=c.line._get_or_add_ln()
    if dashed: ln.append(ln.makeelement(qn("a:prstDash"),{"val":"dash"}))
    if arrow: ln.append(ln.makeelement(qn("a:tailEnd"),{"type":"triangle","w":"med","len":"med"}))
    return c
def hdr(s,cx,t): label(s,cx-200,150,400,28,[P(t,11,FAINT,b=True,align=PP_ALIGN.CENTER,font=HEAD)])

# ============ SLIDE 1: simplified architecture ============
def architecture():
    s=slide()
    label(s,70,34,1780,60,[P("Two Agents, One Trustworthy Spine",32,FG,b=True,font=HEAD)])
    label(s,70,92,1780,30,[P("Two agents share one deterministic, guardrailed spine.",15,MUTE)])
    hdr(s,185,"REQUESTS"); hdr(s,550,"ORCHESTRATION"); hdr(s,1015,"SHARED TRUSTWORTHY SPINE"); hdr(s,1560,"DATA SOURCES")
    box(s,70,210,230,120,line=YEL,lw=2,paras=[P("Trip Planner",15,YEL_LT,b=True,font=HEAD,align=PP_ALIGN.CENTER),P("“where to fish?”",11,MUTE,align=PP_ALIGN.CENTER,sb=2)])
    box(s,70,470,230,120,line=LAV,lw=2,paras=[P("Prospector",15,LAV_LT,b=True,font=HEAD,align=PP_ALIGN.CENTER),P("“find new trout water?”",11,MUTE,align=PP_ALIGN.CENTER,sb=2)])
    box(s,360,210,420,150,line=YEL,lw=2,paras=[P("Trip Planner",18,YEL_LT,b=True,font=HEAD),P("hand-written tool loop",11,MUTE,sb=2),P("Haiku retrieval → Sonnet ranking",11,FG,font=MONO,sb=4)])
    box(s,360,470,420,150,line=LAV,lw=2,paras=[P("Prospector",18,LAV_LT,b=True,font=HEAD),P("LangGraph",11,MUTE,sb=2),P("branching · human-in-the-loop",11,FG,sb=4)])
    box(s,800,180,430,690,fill=PANEL_DK,line=YEL,lw=2,radius=0.04,anchor=MSO_ANCHOR.TOP)
    def stage(y,h,color,title,tag=None,emph=False):
        box(s,820,y,390,h,fill=(PANEL2 if emph else PANEL),line=color,lw=(3 if emph else 2),anchor=MSO_ANCHOR.MIDDLE,paras=[P(title,16,FG,b=True,font=HEAD)])
        if tag: label(s,1030,y+(h/2)-11,168,22,[P(tag,10,color,b=True,align=PP_ALIGN.RIGHT)])
    stage(220,82,YEL,"MCP tool belt",tag="retrieval")
    stage(340,92,GREEN,"Deterministic scorer",tag="the oracle",emph=True)
    stage(470,82,YEL,"Grounding contract")
    stage(580,92,RED,"Guardrail veto",tag="rules decide")
    label(s,820,710,390,26,[P("the model advises · the rules decide",13,YEL_LT,b=True,align=PP_ALIGN.CENTER,font=HEAD)])
    conn(s,1015,302,1015,340,color=MUTE); conn(s,1015,432,1015,470,color=MUTE); conn(s,1015,552,1015,580,color=MUTE)
    for i,nm in enumerate(["USGS NWIS","USGS NLDI","NOAA","State ArcGIS","PAD-US","Postgres"]):
        gx=[1280,1570]; gy=[210,300,390]
        box(s,gx[i%2],gy[i//2],250,72,line=LINE,lw=1.4,anchor=MSO_ANCHOR.MIDDLE,paras=[P(nm,14,FG,b=True,font=HEAD,align=PP_ALIGN.CENTER)])
    conn(s,1280,246,1212,265,color=YEL,width=2.2); label(s,1140,252,90,18,[P("fetch",10,YEL_LT,align=PP_ALIGN.RIGHT)])
    label(s,1560,478,300,22,[P("OUTPUT",11,FAINT,b=True,align=PP_ALIGN.CENTER,font=HEAD)])
    box(s,1280,506,275,120,line=GREEN,lw=2,anchor=MSO_ANCHOR.MIDDLE,paras=[P("$0.02",34,GREEN,b=True,align=PP_ALIGN.CENTER,font=HEAD),P("per decision",11,MUTE,align=PP_ALIGN.CENTER)])
    box(s,1575,506,275,120,line=YEL,lw=2,anchor=MSO_ANCHOR.MIDDLE,paras=[P("17–18 s",30,YEL_LT,b=True,align=PP_ALIGN.CENTER,font=HEAD),P("latency",11,MUTE,align=PP_ALIGN.CENTER)])
    label(s,1280,640,570,40,[P("→ ranked, grounded recommendations + guardrail verdicts, on the map",10.5,MUTE,align=PP_ALIGN.CENTER)])
    conn(s,300,270,356,270,color=YEL,width=2.2); conn(s,300,530,356,530,color=LAV,width=2.2)
    conn(s,780,285,798,300,color=YEL,width=2.2); conn(s,780,525,798,470,color=LAV,width=2.2)
    conn(s,1230,610,1278,566,color=RED,width=2.4)
    band=box(s,70,900,1780,110,fill=PANEL_DK,line=YEL,lw=1.6,radius=0.04,anchor=MSO_ANCHOR.TOP)
    band.line._get_or_add_ln().append(band.line._get_or_add_ln().makeelement(qn("a:prstDash"),{"val":"dash"}))
    label(s,92,924,1000,24,[P("OFFLINE EVAL HARNESS",11,YEL_LT,b=True,font=HEAD)])
    label(s,92,956,1700,28,[P("The scorer here is the eval oracle — same code · validated on 25 planner scenarios + an honest discovery eval",11,MUTE)])
    conn(s,812,900,816,432,color=GREEN,width=1.8,dashed=True)

# ============ SLIDE 2: v0->v3 staircase ============
def staircase():
    s=slide()
    label(s,70,36,1780,50,[P("Trip Planner: The v0 → v3 Staircase",32,FG,b=True,font=HEAD)])
    label(s,70,96,1780,40,[P("Grounding + guardrails turn a confident liar into a trustworthy assistant (25 scenarios)",15,MUTE)])
    cols=[("v0","naive prompt, no tools",RED),("v1","tool-grounded",LAV),("v2","+ catch-log memory",YEL),("v3","+ guardrails & grounding",GREEN)]
    rows=[("Recommendation agreement","higher is better",["8%","100%","100%","100%"],"high"),
          ("Safety violations","lower is better",["16%","0%","0%","0%"],"low"),
          ("Hallucinated readings","lower is better",["100%","4%","12%","0%"],"low")]
    LBLW=330; x0=70; gap=18; cw=(1850-x0-LBLW-gap*4)/4; top=170; headh=104; rowh=150; rgap=14
    for i,(v,sub,color) in enumerate(cols):
        cx=x0+LBLW+i*(cw+gap); emph=(v=="v3")
        box(s,cx,top,cw,headh,fill=(PANEL2 if emph else PANEL),line=color,lw=(3 if emph else 2),anchor=MSO_ANCHOR.MIDDLE,
            paras=[P(v,30,color,b=True,align=PP_ALIGN.CENTER,font=HEAD),P(sub,11,(FG if emph else MUTE),align=PP_ALIGN.CENTER)])
    ry=top+headh+22
    for ri,(lab,note,vals,good) in enumerate(rows):
        yy=ry+ri*(rowh+rgap)
        label(s,x0,yy,LBLW-10,rowh,[P(lab,17,FG,b=True,font=HEAD),P(note,11,FAINT,sb=2)],anchor=MSO_ANCHOR.MIDDLE)
        for ci,(v,sub,color) in enumerate(cols):
            cx=x0+LBLW+ci*(cw+gap); val=vals[ci]; num=float(val.replace("%",""))
            if good=="high": cell=GREEN if num>=90 else (YEL if num>=50 else RED)
            else: cell=GREEN if num==0 else (YEL if num<=15 else RED)
            box(s,cx,yy,cw,rowh,line=cell,lw=(2.6 if v=="v3" else 1.6),anchor=MSO_ANCHOR.MIDDLE,paras=[P(val,40,cell,b=True,align=PP_ALIGN.CENTER,font=HEAD)])
    label(s,70,1004,1780,30,[P("v2’s hallucination bump (4% → 12%) is real — memory added unsourced numbers; the v3 grounding contract drove it to 0%.",14,MUTE)])

# ============ SLIDE 3: orchestration A/B ============
def orchestration():
    s=slide()
    label(s,70,36,1780,50,[P("Engineering Judgment: Right Tool for the Job",32,FG,b=True,font=HEAD)])
    label(s,70,150,1780,30,[P("SAME v3 PLANNER · 25 SCENARIOS · ONLY ORCHESTRATION CHANGES",13,FAINT,b=True,align=PP_ALIGN.CENTER,font=HEAD)])
    cards=[(360,"Hand-written loop",YEL,YEL_LT,"17","lines of orchestration","linear planner"),
           (1010,"LangGraph",LAV,LAV_LT,"38","lines (2.2×)","branching + HITL")]
    cy=200; ch=420; cw=560
    for x,name,color,lt,lines,lsub,sub in cards:
        box(s,x,cy,cw,ch,line=color,lw=2.4,radius=0.05,anchor=MSO_ANCHOR.TOP)
        label(s,x,cy+24,cw,66,[P(name,23,lt,b=True,align=PP_ALIGN.CENTER,font=HEAD),P(sub,12,MUTE,align=PP_ALIGN.CENTER,sb=4)])
        label(s,x,cy+106,cw,110,[P("100%",46,GREEN,b=True,align=PP_ALIGN.CENTER,font=HEAD),P("scenario quality",12,MUTE,align=PP_ALIGN.CENTER,sb=6)])
        conn(s,x+60,cy+230,x+cw-60,cy+230,color=LINE,width=1.4,arrow=False)
        label(s,x,cy+240,cw,110,[P(lines,46,lt,b=True,align=PP_ALIGN.CENTER,font=MONO),P(lsub,12,MUTE,align=PP_ALIGN.CENTER,sb=6)])
    label(s,920,cy+152,80,60,[P("=",40,MUTE,b=True,align=PP_ALIGN.CENTER,font=HEAD)])
    label(s,920,cy+286,80,40,[P("vs",22,MUTE,b=True,align=PP_ALIGN.CENTER,font=HEAD)])
    box(s,360,664,1210,340,fill=PANEL_DK,line=LINE,lw=1.5,radius=0.04,anchor=MSO_ANCHOR.TOP,paras=[
        P("Decision",16,FG,b=True,font=HEAD,sa=6),
        P("Hand-loop for the linear Trip Planner — less code, fully legible.",13.5,MUTE),
        P("LangGraph for the branching, human-in-the-loop Prospector — where interrupt + durable checkpoints actually earn their cost.",13.5,MUTE,sa=8),
        P("Frameworks are not a quality lever — claiming so would be a confound.",13.5,YEL_LT,b=True)])

architecture(); staircase(); orchestration()
path=os.path.join(OUT,"deck_graphics.pptx"); prs.save(path)
print("wrote",path,os.path.getsize(path),"bytes,",len(prs.slides._sldIdLst),"slides")
